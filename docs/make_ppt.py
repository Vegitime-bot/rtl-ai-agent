"""
AI 기반 RTL 자동 생성 — 경영진 보고용 PPT 생성 스크립트
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── 색상 팔레트 ──────────────────────────────────────────────
DARK_BG   = RGBColor(0x1A, 0x1A, 0x2E)   # 딥 네이비
ACCENT    = RGBColor(0x16, 0x21, 0x3E)   # 미드 네이비
HIGHLIGHT = RGBColor(0x0F, 0x3C, 0x8A)   # 블루 강조
GOLD      = RGBColor(0xE2, 0xB9, 0x4B)   # 골드
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT     = RGBColor(0xD0, 0xD8, 0xEE)   # 연한 블루-그레이
GREEN_OK  = RGBColor(0x2E, 0xCC, 0x71)
ORANGE    = RGBColor(0xF3, 0x9C, 0x12)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

blank_layout = prs.slide_layouts[6]  # 완전 blank


# ── 헬퍼 ────────────────────────────────────────────────────
def add_rect(slide, x, y, w, h, fill_rgb, alpha=None):
    shape = slide.shapes.add_shape(1, x, y, w, h)
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    return shape

def add_textbox(slide, x, y, w, h, text, size, bold=False,
                color=WHITE, align=PP_ALIGN.LEFT, italic=False, wrap=True):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p  = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size   = Pt(size)
    run.font.bold   = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb

def slide_bg(slide, color=DARK_BG):
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, color)

def header_bar(slide, title, subtitle=None):
    """상단 헤더 바"""
    add_rect(slide, 0, 0, SLIDE_W, Inches(1.1), HIGHLIGHT)
    add_textbox(slide, Inches(0.4), Inches(0.12), Inches(11), Inches(0.6),
                title, 24, bold=True, align=PP_ALIGN.LEFT)
    if subtitle:
        add_textbox(slide, Inches(0.4), Inches(0.65), Inches(11), Inches(0.4),
                    subtitle, 13, color=LIGHT, align=PP_ALIGN.LEFT)


# ══════════════════════════════════════════════════════════════
# SLIDE 1 — 표지
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
slide_bg(sl)

# 장식 사각형
add_rect(sl, 0, Inches(2.8), Inches(0.25), Inches(1.9), GOLD)
add_rect(sl, Inches(0.25), Inches(2.8), SLIDE_W - Inches(0.25), Inches(1.9), ACCENT)

# 메인 타이틀
add_textbox(sl, Inches(0.6), Inches(2.85), Inches(11), Inches(0.8),
            "AI 기반 RTL 자동 생성 시스템", 36, bold=True, align=PP_ALIGN.LEFT)
add_textbox(sl, Inches(0.6), Inches(3.55), Inches(10), Inches(0.5),
            "반도체 설계 생산성 혁신을 위한 AI Agent 도입 제안", 17,
            color=LIGHT, align=PP_ALIGN.LEFT)

# 하단 정보
add_rect(sl, 0, Inches(6.7), SLIDE_W, Inches(0.8), HIGHLIGHT)
add_textbox(sl, Inches(0.5), Inches(6.78), Inches(6), Inches(0.4),
            "반도체 설계팀 | 2026. 03", 12, color=LIGHT)
add_textbox(sl, Inches(9), Inches(6.78), Inches(4), Inches(0.4),
            "CONFIDENTIAL", 11, color=GOLD, align=PP_ALIGN.RIGHT)


# ══════════════════════════════════════════════════════════════
# SLIDE 2 — 문제 정의 / 현재 Pain Point
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
slide_bg(sl)
header_bar(sl, "문제 정의 — 현재 RTL 수정 작업의 한계",
           "반복적·수작업 중심의 현행 프로세스")

# 두 컬럼: 현황 vs 문제
add_rect(sl, Inches(0.4), Inches(1.3), Inches(5.7), Inches(5.2), ACCENT)
add_rect(sl, Inches(6.3), Inches(1.3), Inches(6.6), Inches(5.2), ACCENT)

add_textbox(sl, Inches(0.5), Inches(1.35), Inches(5.5), Inches(0.5),
            "📋  현행 프로세스", 15, bold=True, color=GOLD)
add_textbox(sl, Inches(6.4), Inches(1.35), Inches(6.4), Inches(0.5),
            "⚠️  핵심 문제", 15, bold=True, color=ORANGE)

current_flow = [
    "① 스펙 변경 문서 수신",
    "② 엔지니어가 수작업으로 스펙 분석",
    "③ 기존 RTL 코드 탐색 및 영향 범위 파악",
    "④ 직접 RTL 수정 · 코딩",
    "⑤ 코드 리뷰 및 검증 반복",
    "⑥ 완료까지 수일~수주 소요",
]
issues = [
    "⏱  반복 수정에 엔지니어 공수 집중",
    "🔍  변경 영향 범위 파악 누락 위험",
    "📄  스펙과 코드 불일치 오류 다수",
    "🔄  검토·수정 루프 반복으로 지연",
    "💸  숙련 엔지니어 시간 낭비",
    "📈  설계 복잡도 증가 → 문제 심화",
]

for i, txt in enumerate(current_flow):
    add_textbox(sl, Inches(0.55), Inches(1.95 + i*0.57), Inches(5.4), Inches(0.5),
                txt, 13, color=LIGHT)
for i, txt in enumerate(issues):
    add_textbox(sl, Inches(6.45), Inches(1.95 + i*0.57), Inches(6.2), Inches(0.5),
                txt, 13, color=LIGHT)


# ══════════════════════════════════════════════════════════════
# SLIDE 3 — 솔루션 개요 (핵심 컨셉)
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
slide_bg(sl)
header_bar(sl, "솔루션 개요 — AI Agent 기반 자동화",
           "스펙을 이해하고 스스로 RTL을 생성하는 지능형 파이프라인")

# 중앙 컨셉 박스
add_rect(sl, Inches(3.2), Inches(1.4), Inches(6.9), Inches(1.0), HIGHLIGHT)
add_textbox(sl, Inches(3.2), Inches(1.5), Inches(6.9), Inches(0.8),
            "스펙 변경 입력  →  AI Agent  →  검증된 RTL 자동 출력",
            17, bold=True, align=PP_ALIGN.CENTER)

# 3개 핵심 축
cols = [
    (Inches(0.3),  "🧠  스펙 이해",    "알고리즘·uArch 변경 내용을\nAI가 자동으로 분석·해석"),
    (Inches(4.65), "🔗  구조 추론",    "신호 인과관계 그래프를 통해\n변경 영향 범위를 정밀 파악"),
    (Inches(9.0),  "⚙️  코드 생성",    "LLM이 분석 결과를 바탕으로\n올바른 RTL을 자동 생성·검증"),
]
for cx, title, body in cols:
    add_rect(sl, cx, Inches(2.65), Inches(4.0), Inches(2.2), ACCENT)
    add_textbox(sl, cx + Inches(0.1), Inches(2.72), Inches(3.8), Inches(0.55),
                title, 15, bold=True, color=GOLD)
    add_textbox(sl, cx + Inches(0.1), Inches(3.25), Inches(3.8), Inches(1.5),
                body, 13, color=LIGHT)

# 기대 효과 요약
add_rect(sl, Inches(0.3), Inches(5.1), SLIDE_W - Inches(0.6), Inches(1.6), ACCENT)
add_textbox(sl, Inches(0.5), Inches(5.15), Inches(4), Inches(0.4),
            "🎯  기대 효과", 14, bold=True, color=GOLD)
effects = [
    "✅  수일 걸리던 RTL 수정 → 수십 분 이내 완료",
    "✅  스펙 누락·오역 등 인적 오류 대폭 감소",
    "✅  숙련 엔지니어를 고부가 설계 업무에 집중 투입",
    "✅  설계 복잡도 증가에도 일관된 품질 유지",
]
for i, e in enumerate(effects):
    col = 0 if i < 2 else 1
    row = i % 2
    add_textbox(sl, Inches(0.5 + col * 6.5), Inches(5.7 + row * 0.45),
                Inches(6.2), Inches(0.4), e, 12, color=LIGHT)


# ══════════════════════════════════════════════════════════════
# SLIDE 4 — 파이프라인 흐름 (비기술 버전)
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
slide_bg(sl)
header_bar(sl, "시스템 동작 흐름",
           "입력부터 검증된 RTL 출력까지 4단계 자동화 파이프라인")

# 단계 박스 4개 (가로 흐름)
steps = [
    ("STEP 1", "입력 수신",
     "기존 RTL\n변경 스펙(uArch)\n알고리즘 변경 내용"),
    ("STEP 2", "분석 · 이해",
     "RTL 구조 파싱\n신호 관계 그래프 구축\n알고리즘 변경점 추출"),
    ("STEP 3", "AI 생성",
     "스펙 분석 Agent\n변경 계획 수립 Agent\nLLM 기반 RTL 코드 생성"),
    ("STEP 4", "검증 · 출력",
     "신호 구조 자동 검증\n불일치 시 자동 재생성\n최종 RTL 파일 출력"),
]

for i, (step_lbl, title, body) in enumerate(steps):
    bx = Inches(0.25 + i * 3.27)
    # 상단 라벨
    add_rect(sl, bx, Inches(1.35), Inches(3.0), Inches(0.45), HIGHLIGHT)
    add_textbox(sl, bx, Inches(1.38), Inches(3.0), Inches(0.4),
                step_lbl, 11, bold=True, align=PP_ALIGN.CENTER)
    # 본문 박스
    add_rect(sl, bx, Inches(1.8), Inches(3.0), Inches(3.6), ACCENT)
    add_textbox(sl, bx + Inches(0.1), Inches(1.88), Inches(2.8), Inches(0.55),
                title, 15, bold=True, color=GOLD)
    add_textbox(sl, bx + Inches(0.1), Inches(2.5), Inches(2.8), Inches(2.7),
                body, 13, color=LIGHT)
    # 화살표 (마지막 제외)
    if i < 3:
        add_textbox(sl, bx + Inches(3.0), Inches(2.7), Inches(0.3), Inches(0.6),
                    "▶", 20, color=GOLD, align=PP_ALIGN.CENTER)

# 하단 요약
add_rect(sl, Inches(0.25), Inches(5.6), SLIDE_W - Inches(0.5), Inches(1.2), ACCENT)
add_textbox(sl, Inches(0.5), Inches(5.65), Inches(12), Inches(1.05),
            "기존: 엔지니어가 스펙 해석 → 수동 코딩 → 리뷰 반복 (수일 소요)\n"
            "신규: AI가 스펙을 직접 이해 → 자동 코드 생성 → 자동 검증 (수십 분 이내)",
            13, color=LIGHT)


# ══════════════════════════════════════════════════════════════
# SLIDE 5 — 핵심 기술 차별점 (비기술 설명)
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
slide_bg(sl)
header_bar(sl, "핵심 기술 차별점",
           "단순 코드 생성이 아닌 '설계 의도를 이해하는' AI")

differentiators = [
    ("🔗  신호 인과관계 지식 그래프",
     "RTL 내 모든 신호의 연결 관계를 그래프 DB로 구축.\n"
     "변경 시 영향 받는 신호를 자동으로 추적해 LLM에 전달.\n"
     "단순 텍스트 이해를 넘어 설계 구조를 기반으로 코드 생성."),
    ("📚  스펙 의미 검색 (Vector RAG)",
     "uArch 스펙 문서를 의미 단위로 분할·인덱싱.\n"
     "변경 내용과 관련된 스펙 조각을 자동으로 검색·주입.\n"
     "방대한 문서 중 필요한 맥락만 정확히 LLM에 전달."),
    ("🔄  검증 실패 시 자동 재생성 루프",
     "생성된 RTL의 신호 구조가 스펙과 일치하는지 자동 검증.\n"
     "불일치 발생 시 실패 원인을 AI에게 피드백 → 자동 재시도.\n"
     "사람 개입 없이 합격 기준을 만족할 때까지 반복."),
    ("📐  변경 영향 범위 기반 최적화",
     "전체 RTL이 아닌 변경에 영향 받는 부분만 선택적 처리.\n"
     "대형 RTL에서도 LLM 컨텍스트 한계 없이 안정적 동작.\n"
     "정밀 타겟팅으로 불필요한 코드 변경 최소화."),
]

for i, (title, body) in enumerate(differentiators):
    row = i // 2
    col = i %  2
    bx = Inches(0.3 + col * 6.55)
    by = Inches(1.35 + row * 2.6)
    add_rect(sl, bx, by, Inches(6.3), Inches(2.4), ACCENT)
    add_textbox(sl, bx + Inches(0.15), by + Inches(0.1), Inches(6.0), Inches(0.55),
                title, 14, bold=True, color=GOLD)
    add_textbox(sl, bx + Inches(0.15), by + Inches(0.65), Inches(6.0), Inches(1.6),
                body, 12, color=LIGHT)


# ══════════════════════════════════════════════════════════════
# SLIDE 6 — 현재 진행 상황 및 향후 계획
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
slide_bg(sl)
header_bar(sl, "현재 상태 및 향후 계획",
           "핵심 파이프라인 완성 — 고도화 단계 진입")

# 현황 타임라인
add_textbox(sl, Inches(0.4), Inches(1.25), Inches(5), Inches(0.4),
            "📍 현재 완료된 항목", 14, bold=True, color=GOLD)

done_items = [
    "✅  RTL 파싱 및 구조 분석",
    "✅  신호 인과관계 그래프 (Neo4j)",
    "✅  알고리즘 변경 자동 분석",
    "✅  스펙 의미 검색 (Vector RAG)",
    "✅  LLM 기반 RTL 코드 생성",
    "✅  인과관계 기반 자동 검증 루프",
    "✅  대형 RTL 대응 (청크 최적화)",
    "✅  멀티 파일 입력 지원",
]
for i, item in enumerate(done_items):
    col = i // 4
    row = i %  4
    add_textbox(sl, Inches(0.45 + col * 3.0), Inches(1.75 + row * 0.52),
                Inches(2.85), Inches(0.45), item, 12, color=GREEN_OK)

# 구분선
add_rect(sl, Inches(6.4), Inches(1.25), Inches(0.04), Inches(4.8), HIGHLIGHT)

# 향후 계획
add_textbox(sl, Inches(6.6), Inches(1.25), Inches(5), Inches(0.4),
            "🚀 향후 고도화 방향", 14, bold=True, color=GOLD)

next_items = [
    ("Phase 1\n(단기)", "시뮬레이션 연동\n— 기능 검증 자동화 추가"),
    ("Phase 2\n(중기)", "다국어 스펙 지원\n— 자연어 스펙 직접 입력"),
    ("Phase 3\n(장기)", "설계 플랫폼 통합\n— EDA 툴체인 직접 연동"),
]
for i, (phase, body) in enumerate(next_items):
    bx = Inches(6.6)
    by = Inches(1.85 + i * 1.35)
    add_rect(sl, bx, by, Inches(0.8), Inches(1.1), HIGHLIGHT)
    add_textbox(sl, bx, by + Inches(0.1), Inches(0.8), Inches(0.9),
                phase, 10, bold=True, align=PP_ALIGN.CENTER, color=WHITE)
    add_textbox(sl, Inches(7.55), by + Inches(0.05), Inches(5.2), Inches(1.0),
                body, 13, color=LIGHT)


# ══════════════════════════════════════════════════════════════
# SLIDE 7 — 기대 효과 / 비즈니스 임팩트
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
slide_bg(sl)
header_bar(sl, "기대 효과 및 비즈니스 임팩트",
           "생산성·품질·경쟁력 세 축의 동시 향상")

metrics = [
    ("⏱", "설계 사이클\n단축", "수일 → 수십 분", "RTL 수정 리드타임 대폭 감소"),
    ("🔧", "인적 오류\n감소", "스펙 누락·오역\n자동 차단", "검증 루프로 품질 보증"),
    ("👷", "엔지니어\n역할 변화", "반복 코딩\n→ 고부가 설계", "핵심 인력 효율 극대화"),
    ("📈", "설계 복잡도\n대응력", "규모 증가에도\n일관된 성능", "경쟁사 대비 민첩성 확보"),
]

for i, (icon, title, stat, desc) in enumerate(metrics):
    bx = Inches(0.3 + i * 3.27)
    add_rect(sl, bx, Inches(1.35), Inches(3.0), Inches(4.8), ACCENT)
    add_textbox(sl, bx, Inches(1.5),  Inches(3.0), Inches(0.7),
                icon, 28, align=PP_ALIGN.CENTER)
    add_textbox(sl, bx, Inches(2.15), Inches(3.0), Inches(0.65),
                title, 14, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    add_textbox(sl, bx, Inches(2.85), Inches(3.0), Inches(0.8),
                stat, 13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(sl, bx + Inches(0.1), Inches(3.7), Inches(2.8), Inches(1.0),
                desc, 12, color=LIGHT, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════
# SLIDE 8 — 마무리 / 요청 사항
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
slide_bg(sl)

add_rect(sl, 0, Inches(2.5), Inches(0.3), Inches(2.5), GOLD)
add_rect(sl, Inches(0.3), Inches(2.5), SLIDE_W - Inches(0.3), Inches(2.5), ACCENT)

add_textbox(sl, Inches(0.7), Inches(2.6), Inches(11), Inches(0.8),
            "요청 사항", 30, bold=True)
add_textbox(sl, Inches(0.7), Inches(3.3), Inches(11.5), Inches(1.5),
            "본 AI Agent 시스템의 사내 도입 및 고도화를 위한\n"
            "추가 개발 리소스 및 인프라 지원을 요청드립니다.\n"
            "단계적 도입을 통해 빠른 ROI 실현이 가능합니다.",
            15, color=LIGHT)

add_rect(sl, 0, Inches(6.7), SLIDE_W, Inches(0.8), HIGHLIGHT)
add_textbox(sl, Inches(0.5), Inches(6.78), Inches(6), Inches(0.4),
            "반도체 설계팀 | 2026. 03", 12, color=LIGHT)
add_textbox(sl, Inches(8.5), Inches(6.78), Inches(4.5), Inches(0.4),
            "AI 기반 RTL 자동 생성 시스템", 12, color=GOLD, align=PP_ALIGN.RIGHT)


# ── 저장 ────────────────────────────────────────────────────
out_path = "/Users/vegitime/.openclaw/workspace/projects/rtl-ai-agent/docs/RTL_AI_Agent_경영진보고.pptx"
prs.save(out_path)
print(f"✅ PPT 저장 완료: {out_path}")
