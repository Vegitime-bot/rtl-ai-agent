"""
기존 PPT에 다이어그램 슬라이드 2장 추가
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import copy

DARK_BG   = RGBColor(0x1A, 0x1A, 0x2E)
HIGHLIGHT = RGBColor(0x0F, 0x3C, 0x8A)
GOLD      = RGBColor(0xE2, 0xB9, 0x4B)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT     = RGBColor(0xD0, 0xD8, 0xEE)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

PPT_PATH  = "/Users/vegitime/.openclaw/workspace/projects/rtl-ai-agent/docs/RTL_AI_Agent_경영진보고.pptx"
IMG1_PATH = "/Users/vegitime/.openclaw/workspace/projects/rtl-ai-agent/docs/diagram_pipeline.png"
IMG2_PATH = "/Users/vegitime/.openclaw/workspace/projects/rtl-ai-agent/docs/diagram_concept.png"

prs = Presentation(PPT_PATH)
blank = prs.slide_layouts[6]

def add_rect(slide, x, y, w, h, fill_rgb):
    from pptx.util import Emu
    shape = slide.shapes.add_shape(1, x, y, w, h)
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    return shape

def add_textbox(slide, x, y, w, h, text, size, bold=False, color=WHITE, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tb

def make_diagram_slide(prs, title, subtitle, img_path, insert_at):
    # 슬라이드 3 뒤에 삽입
    xml_slides = prs.slides._sldIdLst
    sl = prs.slides.add_slide(blank)
    # 배경
    add_rect(sl, 0, 0, SLIDE_W, SLIDE_H, DARK_BG)
    # 헤더
    add_rect(sl, 0, 0, SLIDE_W, Inches(1.0), HIGHLIGHT)
    add_textbox(sl, Inches(0.4), Inches(0.1), Inches(12), Inches(0.55),
                title, 22, bold=True)
    add_textbox(sl, Inches(0.4), Inches(0.6), Inches(12), Inches(0.35),
                subtitle, 12, color=LIGHT)
    # 이미지 (헤더 아래 전체)
    sl.shapes.add_picture(img_path, Inches(0.2), Inches(1.05),
                          SLIDE_W - Inches(0.4), SLIDE_H - Inches(1.15))
    # 하단 바
    add_rect(sl, 0, Inches(7.2), SLIDE_W, Inches(0.3), HIGHLIGHT)
    return sl

# SLIDE 4 위치 (인덱스 3) 뒤에 삽입 → 슬라이드 4, 5번 위치
make_diagram_slide(
    prs,
    "파이프라인 흐름 — 단계별 자동화",
    "4단계: 입력 수신 → 분석·빌드 → AI 추론·생성 → 검증·출력",
    IMG1_PATH,
    insert_at=3
)
make_diagram_slide(
    prs,
    "핵심 기술 컨셉 — AI Agent 아키텍처",
    "신호 인과관계 그래프 × Vector RAG × LLM 오케스트레이션",
    IMG2_PATH,
    insert_at=4
)

# 저장
prs.save(PPT_PATH)
print(f"✅ PPT 업데이트 완료 (총 {len(prs.slides)}슬라이드): {PPT_PATH}")
